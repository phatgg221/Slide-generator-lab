"""Skill 4: write generated content back into the .pptx template.

Text replacement keeps the template's design: each paragraph keeps its first
run's font/color/size and extra runs are dropped; multi-line content maps one
line per paragraph, cloning the last paragraph's XML when more lines are
needed so bullets keep their styling.

Image replacement adds the new picture as a fresh image part and repoints the
existing picture shape's blip relationship to it, so position, size, crop,
and frame effects from the template survive untouched.
"""

from __future__ import annotations

import copy
import io
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from .content_generator import GeneratedDeckContent
from .template_parser import TemplateSpec, iter_leaf_shapes

_MIN_FONT_SCALE = 0.55   # never shrink below 55% of the designed size


def _set_paragraph_text(paragraph, text: str) -> None:
    """Put text in a paragraph, keeping the first run's formatting."""
    runs = paragraph.runs
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        paragraph.add_run().text = text


def set_text_preserving_format(text_frame, new_text: str) -> None:
    """Replace a text frame's content line-by-line. Lines map onto existing
    paragraphs; surplus lines clone the last paragraph's XML (so bullet
    styling carries over), surplus paragraphs are removed."""
    lines = new_text.split("\n") or [""]
    paragraphs = list(text_frame.paragraphs)

    while len(paragraphs) < len(lines):
        clone = copy.deepcopy(paragraphs[-1]._p)
        paragraphs[-1]._p.addnext(clone)
        paragraphs = list(text_frame.paragraphs)

    for para, line in zip(paragraphs, lines):
        _set_paragraph_text(para, line)

    for para in paragraphs[len(lines):]:
        para._p.getparent().remove(para._p)


def shrink_to_fit(text_frame, new_len: int, budget: int) -> float:
    """When text exceeds its budget, scale every explicit font size down so
    the longer text occupies roughly the designed area (text fills space in
    two dimensions, hence the square root). Returns the factor applied."""
    if new_len <= budget * 1.05:
        return 1.0
    factor = max(_MIN_FONT_SCALE, (budget / new_len) ** 0.5)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                run.font.size = Pt(round(run.font.size.pt * factor, 1))
    return factor


def replace_picture_image(picture, image_bytes: bytes) -> None:
    """Swap the bitmap behind a picture shape without touching its geometry,
    crop, or effects."""
    image_part, rId = picture.part.get_or_add_image_part(io.BytesIO(image_bytes))
    picture._element.blipFill.blip.rEmbed = rId


def fill_template(
    template_path: Union[str, Path],
    content: GeneratedDeckContent,
    output_path: Union[str, Path],
    images: dict[tuple[int, int], bytes] | None = None,
    spec: TemplateSpec | None = None,
) -> Path:
    """Apply generated content to a template and save the result.

    Args:
        template_path: the Canva-exported .pptx.
        content: shape_id-keyed text from the content_generator skill.
        output_path: where to save the filled deck.
        images: optional {(slide_index, shape_id): png_bytes} from the
            image_generator skill.
        spec: when given, text exceeding its budget gets its font shrunk
            proportionally so it stays inside the designed box.

    Returns the output path.
    """
    images = images or {}
    budgets = {}
    if spec is not None:
        budgets = {(s.index, t.shape_id): t.max_chars
                   for s in spec.slides for t in s.texts}
    prs = Presentation(str(template_path))

    for idx, slide in enumerate(prs.slides):
        slide_content = content.slide(idx)
        for shape in iter_leaf_shapes(slide.shapes):
            if slide_content and shape.has_text_frame and shape.shape_id in slide_content.texts:
                text = slide_content.texts[shape.shape_id]
                set_text_preserving_format(shape.text_frame, text)
                budget = budgets.get((idx, shape.shape_id))
                if budget:
                    shrink_to_fit(shape.text_frame, len(text), budget)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and (idx, shape.shape_id) in images:
                replace_picture_image(shape, images[(idx, shape.shape_id)])

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
