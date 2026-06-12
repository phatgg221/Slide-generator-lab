"""Skill 12: element entrance animations (PowerPoint-native).

Writes <p:timing> XML so shapes appear one after another when presenting —
the programmatic equivalent of PowerPoint's Fade/Wipe/Appear entrance
effects. This animates the deck itself (unlike Canva, whose animations never
reach the .pptx).

    add_animations("deck.pptx", "deck.pptx", effect="fade", trigger="auto")

Scope: top-level shapes only (animating inside groups is unreliable across
PowerPoint versions), entrance effects only. trigger="auto" plays the
sequence after the slide appears; trigger="click" advances per click.
"""

from __future__ import annotations

from pathlib import Path
from typing import Union

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# effect -> (presetID, animEffect filter or None for plain appear)
ENTRANCE_EFFECTS = {
    "appear": (1, None),
    "fade":   (10, "fade"),
    "wipe":   (22, "wipe(up)"),
    "blinds": (3, "blinds(horizontal)"),
}


def _xml(fragment: str):
    return etree.fromstring(
        f'<root xmlns:p="{_P}">{fragment}</root>')[0]


def _effect_xml(ids, spid, preset_id, filt, duration_ms, delay_ms, node_type):
    """One entrance effect (innermost p:par with the preset)."""
    anim = ""
    if filt:
        anim = (
            f'<p:animEffect transition="in" filter="{filt}">'
            f'<p:cBhvr><p:cTn id="{ids[2]}" dur="{duration_ms}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'</p:animEffect>'
        )
    return (
        f'<p:par><p:cTn id="{ids[0]}" presetID="{preset_id}" presetClass="entr" '
        f'presetSubtype="0" fill="hold" grpId="0" nodeType="{node_type}">'
        f'<p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{ids[1]}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'{anim}'
        f'</p:childTnLst></p:cTn></p:par>'
    )


def _click_group(ids, effect_xml):
    """A click-triggered group holding one effect (PowerPoint's per-click
    structure: outer par waits indefinitely for the click)."""
    return _xml(
        f'<p:par><p:cTn id="{ids[0]}" fill="hold">'
        f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
        f'<p:par><p:cTn id="{ids[1]}" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'{effect_xml}'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
    )


def _auto_group(ids, effects_xml):
    """One auto-starting group holding ALL effects as an after-previous
    chain. The dual condition (indefinite + onBegin of the main sequence,
    tn val=2) is PowerPoint's native marker for 'starts with the slide'."""
    return _xml(
        f'<p:par><p:cTn id="{ids[0]}" fill="hold">'
        f'<p:stCondLst><p:cond delay="indefinite"/>'
        f'<p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:par><p:cTn id="{ids[1]}" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'{effects_xml}'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
    )


def _animatable_shapes(slide, include_pictures):
    shapes = []
    for shape in slide.shapes:  # top level only, by design
        has_text = shape.has_text_frame and shape.text_frame.text.strip()
        is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        if has_text or is_group or (is_pic and include_pictures):
            shapes.append(shape)
    shapes.sort(key=lambda s: ((s.top or 0), (s.left or 0)))
    return shapes


def add_animations(
    pptx_path: Union[str, Path],
    output_path: Union[str, Path],
    *,
    effect: str = "fade",
    trigger: str = "auto",        # "auto" (after slide appears) | "click"
    duration_ms: int = 500,
    stagger_ms: int = 400,
    include_pictures: bool = True,
    max_per_slide: int = 12,
) -> Path:
    """Add an entrance animation to every text/picture/group shape, playing
    top-to-bottom. Replaces any existing animations."""
    if effect not in ENTRANCE_EFFECTS:
        raise ValueError(f"Unknown effect {effect!r}. Available: {sorted(ENTRANCE_EFFECTS)}")
    if trigger not in ("auto", "click"):
        raise ValueError("trigger must be 'auto' or 'click'")
    preset_id, filt = ENTRANCE_EFFECTS[effect]

    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        shapes = _animatable_shapes(slide, include_pictures)[:max_per_slide]
        sld = slide._element
        for old in sld.findall(qn("p:timing")):
            sld.remove(old)
        if not shapes:
            continue

        next_id = 5
        group_pars, bld_entries = [], []
        if trigger == "auto":
            chunks = []
            for i, shape in enumerate(shapes):
                ids = list(range(next_id, next_id + 3))
                next_id += 3
                chunks.append(_effect_xml(
                    ids, shape.shape_id, preset_id, filt, duration_ms,
                    i * stagger_ms, "afterEffect"))
                bld_entries.append(f'<p:bldP spid="{shape.shape_id}" grpId="0"/>')
            group_pars.append(_auto_group([3, 4], "".join(chunks)))
        else:
            for shape in shapes:
                ids = list(range(next_id, next_id + 5))
                next_id += 5
                effect = _effect_xml(
                    ids[2:], shape.shape_id, preset_id, filt, duration_ms,
                    0, "clickEffect")
                group_pars.append(_click_group(ids[:2], effect))
                bld_entries.append(f'<p:bldP spid="{shape.shape_id}" grpId="0"/>')

        timing = _xml(
            '<p:timing><p:tnLst><p:par>'
            '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
            '<p:childTnLst></p:childTnLst></p:cTn>'
            '<p:prevCondLst><p:cond evt="onPrev" delay="0">'
            '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            '<p:nextCondLst><p:cond evt="onNext" delay="0">'
            '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
            f'<p:bldLst>{"".join(bld_entries)}</p:bldLst></p:timing>'
        )
        main_seq_children = timing.find(
            f'.//{{{_P}}}cTn[@nodeType="mainSeq"]/{{{_P}}}childTnLst')
        for par in group_pars:
            main_seq_children.append(par)
        sld.append(timing)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
