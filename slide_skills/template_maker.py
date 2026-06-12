"""Skill 10: turn any uploaded .pptx into a registered, fillable template.

This is the ingestion step that runs BEFORE generation:

    user's old slides -> clean_template -> auto_manifest (GPT-4o) -> library

    clean_template    -- delete Canva tip bubbles and dead navigation buttons
                         (visible junk in generated decks), drop unwanted slides
    auto_manifest     -- GPT-4o looks at each slide's content and assigns a
                         type name + description (marks unusable ones _skip_*)
    prepare_template  -- the full pipeline: clean + classify + save into the
                         library folder as <name>.pptx + <name>.json
    list_templates    -- the registry: every template available for generation
                         (for the FastAPI app / chatbot to offer as choices)
"""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import get_client, TEXT_MODEL
from .template_parser import (
    DEFAULT_EXCLUDE_PREFIXES, _is_navigation, parse_template,
)
from .usage import tracker

DEFAULT_LIBRARY_DIR = Path("library")


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _clean_shapes(shapes, tip_prefixes) -> int:
    """Recursively remove tip bubbles and navigation buttons. Returns count."""
    removed = 0
    for shape in list(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_texts = [
                s.text_frame.text for s in shape.shapes
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            # a group whose text is only tips is a tip bubble with decoration
            if group_texts and all(t.strip().startswith(tip_prefixes) for t in group_texts):
                _delete_shape(shape)
                removed += 1
            else:
                removed += _clean_shapes(shape.shapes, tip_prefixes)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            if text.startswith(tip_prefixes) or _is_navigation(shape.text_frame):
                _delete_shape(shape)
                removed += 1
    return removed


def clean_template(
    src: Union[str, Path],
    dst: Union[str, Path],
    *,
    tip_prefixes: tuple[str, ...] = DEFAULT_EXCLUDE_PREFIXES,
    drop_slides: list[int] | None = None,
) -> Path:
    """Copy src to dst with junk removed: tip bubbles, navigation buttons,
    and the slides whose indices are in drop_slides."""
    prs = Presentation(str(src))

    if drop_slides:
        sld_id_lst = prs.slides._sldIdLst
        for i, sld_id in reversed(list(enumerate(list(sld_id_lst)))):
            if i in set(drop_slides):
                r_id = sld_id.rId
                sld_id_lst.remove(sld_id)
                prs.part.drop_rel(r_id)

    removed = 0
    for slide in prs.slides:
        removed += _clean_shapes(slide.shapes, tip_prefixes)

    dst = Path(dst)
    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    return dst


def auto_manifest(pptx_path: Union[str, Path]) -> dict:
    """GPT-4o reads each slide's content and produces the library manifest:
    a unique snake_case type per slide plus a planner-facing description.
    Slides with nothing fillable get _skip_* types."""
    spec = parse_template(pptx_path)
    slides_summary = []
    for s in spec.slides:
        slides_summary.append({
            "index": s.index,
            "texts": [
                {"role": t.role, "sample": t.current_text[:80],
                 "max_chars": t.max_chars, "paragraphs": t.n_paragraphs}
                for t in s.texts
            ],
            "image_placeholders": len(s.images),
        })

    client = get_client()
    response = client.chat.completions.create(
        model=TEXT_MODEL,
        temperature=0.2,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": (
                "You are classifying the slides of a presentation template so "
                "an automated planner can reuse them. For EVERY slide, assign:\n"
                "- type: short unique snake_case name describing the slide's "
                "function (e.g. title, agenda, section_divider, concept, "
                "bullets, statistic, comparison, quote, gallery, process, "
                "summary, closing). Same design used twice -> suffix _2.\n"
                "- description: one sentence telling the planner when to use "
                "it and what content it holds (mention text box count and "
                "image placeholders).\n"
                "If a slide has no usable content boxes (tutorial pages, "
                "decorative filler), name it _skip_<reason>.\n"
                'Return ONLY JSON: {"slides": [{"index": 0, "type": "...", '
                '"description": "..."}]} covering every index.'
            )},
            {"role": "user", "content": json.dumps(slides_summary, ensure_ascii=False)},
        ],
    )
    tracker.record_chat(response.usage)
    payload = json.loads(response.choices[0].message.content)

    by_index = {s.get("index"): s for s in payload.get("slides", [])}
    types, descriptions, seen = [], {}, set()
    for s in spec.slides:
        entry = by_index.get(s.index, {})
        t = re.sub(r"[^a-z0-9_]", "_", str(entry.get("type", f"slide_{s.index}")).lower())
        if not s.texts and not s.images and not t.startswith("_"):
            t = f"_skip_{t}"
        base = t
        n = 2
        while t in seen:
            t = f"{base}_{n}"
            n += 1
        seen.add(t)
        types.append(t)
        if not t.startswith("_"):
            descriptions[t] = str(entry.get("description", ""))
    return {"types": types, "descriptions": descriptions}


def prepare_template(
    src: Union[str, Path],
    name: str,
    *,
    library_dir: Union[str, Path] = DEFAULT_LIBRARY_DIR,
    drop_slides: list[int] | None = None,
    classify: bool = True,
) -> dict:
    """Full ingestion: clean the deck, classify its slides, and register it
    in the library. Returns {"pptx", "manifest", "types"}."""
    name = re.sub(r"[^A-Za-z0-9_-]", "_", name)
    library_dir = Path(library_dir)
    pptx_path = library_dir / f"{name}.pptx"
    manifest_path = library_dir / f"{name}.json"

    clean_template(src, pptx_path, drop_slides=drop_slides)

    if classify:
        manifest = auto_manifest(pptx_path)
    else:
        n = len(Presentation(str(pptx_path)).slides._sldIdLst)
        manifest = {"types": [f"slide_{i}" for i in range(n)], "descriptions": {}}

    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"pptx": str(pptx_path), "manifest": str(manifest_path),
            "types": manifest["types"]}


def list_templates(library_dir: Union[str, Path] = DEFAULT_LIBRARY_DIR) -> list[dict]:
    """Registry of ready-to-use templates (for an API/chatbot to offer)."""
    out = []
    for pptx in sorted(Path(library_dir).glob("*.pptx")):
        manifest = pptx.with_suffix(".json")
        if not manifest.exists():
            continue
        data = json.loads(manifest.read_text(encoding="utf-8"))
        usable = [t for t in data["types"] if not t.startswith("_")]
        out.append({
            "name": pptx.stem,
            "path": str(pptx),
            "slide_types": usable,
            "descriptions": data.get("descriptions", {}),
        })
    return out
