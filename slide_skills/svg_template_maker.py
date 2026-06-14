"""Skill 14: convert a .pptx template into a folder of SVG templates.

The output SVGs have the original design as a pixel-perfect vector
background and live {{placeholder}} <text> elements where the fillable text
was — the "SVG template" format that Canva can't export.

Pipeline:
    1. copy the deck with every fillable text blanked (design stays)
    2. render that copy to PDF (Microsoft PowerPoint via AppleScript on
       macOS, or LibreOffice if available)
    3. each PDF page -> SVG background (PyMuPDF, vector)
    4. inject {{name}} <text> overlays at the original geometry
       (position, size, color, weight, alignment from the .pptx)

Output folder:
    <out>/slide_00_title.svg, ... + schema.json (placeholder metadata)

Render filled slides with svg_slide_renderer.render_svg_deck.
"""

from __future__ import annotations

import base64
import html
import json
import re
import shutil
import subprocess
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from .template_parser import iter_leaf_shapes, parse_template
from .slide_filler import set_text_preserving_format

_APPLESCRIPT = '''
tell application "Microsoft PowerPoint"
    activate
    set thePres to open (POSIX file "{src}")
    delay 1
    save thePres in (POSIX file "{dst}") as save as PDF
    delay 1
    close thePres saving no
end tell
'''


def _blank_fillable(src: Path, dst: Path, spec) -> None:
    prs = Presentation(str(src))
    slides = list(prs.slides)
    for s in spec.slides:
        fillable = {t.shape_id for t in s.texts}
        for shape in iter_leaf_shapes(slides[s.index].shapes):
            if shape.shape_id in fillable and shape.has_text_frame:
                set_text_preserving_format(shape.text_frame, "")
    prs.save(str(dst))


def _pptx_to_pdf(src: Path, dst: Path) -> None:
    soffice = shutil.which("soffice")
    if soffice:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(dst.parent), str(src)],
            check=True, capture_output=True, timeout=300)
        produced = dst.parent / (src.stem + ".pdf")
        if produced != dst:
            produced.rename(dst)
        return
    script = _APPLESCRIPT.format(src=src, dst=dst)
    result = subprocess.run(["osascript", "-e", script],
                            capture_output=True, text=True, timeout=300)
    if result.returncode != 0 or not dst.exists():
        raise RuntimeError(
            "PDF export failed. PowerPoint automation said: "
            f"{result.stderr.strip() or result.stdout.strip()}\n"
            "(macOS may have shown a permission dialog — approve it and retry, "
            "or install LibreOffice.)")


def _text_color(element) -> str:
    return "#000000"  # overridden below when a run color is found


def _shape_text_style(shape):
    """(font_pt, color, bold, align) from the first non-empty run."""
    font_pt, color, bold, align = None, "#000000", False, "left"
    for para in shape.text_frame.paragraphs:
        if para.alignment == PP_ALIGN.CENTER:
            align = "middle"
        elif para.alignment == PP_ALIGN.RIGHT:
            align = "end"
        for run in para.runs:
            if run.font.size is not None:
                font_pt = run.font.size.pt
            try:
                if run.font.color and run.font.color.rgb is not None:
                    color = f"#{run.font.color.rgb}"
            except (AttributeError, TypeError):
                pass
            bold = bool(run.font.bold)
            return font_pt, color, bold, align
    return font_pt, color, bold, align


def make_svg_templates(
    pptx_path: Union[str, Path],
    out_dir: Union[str, Path],
    *,
    types: list[str] | None = None,
    font_family: str = "Helvetica, Arial, sans-serif",
) -> dict:
    """Convert each slide of pptx_path into <out_dir>/slide_NN_<type>.svg
    with {{placeholders}}, plus schema.json. Returns the schema."""
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    if types is None:
        manifest = pptx_path.with_suffix(".json")
        if manifest.exists():
            types = json.loads(manifest.read_text(encoding="utf-8")).get("types")

    spec = parse_template(pptx_path)
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    # pt per inch is 72; SVG coords below are in points to match the PDF
    page_w = prs.slide_width / 914400 * 72
    page_h = prs.slide_height / 914400 * 72

    # Work files live next to the output: sandboxed PowerPoint silently fails
    # to open files under /tmp, then 'active presentation' grabs whatever the
    # user has open instead.
    work = out_dir / ".work"
    work.mkdir(exist_ok=True)
    if True:
        blanked = work / f"blanked_{pptx_path.stem}.pptx"
        pdf = work / f"blanked_{pptx_path.stem}.pdf"
        _blank_fillable(pptx_path, blanked, spec)
        _pptx_to_pdf(blanked, pdf)

        import fitz
        doc = fitz.open(str(pdf))
        schema: dict = {"slides": []}

        for s in spec.slides:
            slide_type = (types[s.index] if types and s.index < len(types)
                          else f"slide_{s.index}")
            page = doc[s.index]
            # background: render the vector page to high-res PNG and embed it
            # (PyMuPDF SVG output embeds fonts python renderers choke on,
            # so a 2x raster background is the robust choice)
            pix = page.get_pixmap(matrix=fitz.Matrix(2.5, 2.5))
            b64 = base64.b64encode(pix.tobytes("png")).decode()

            overlays, slide_schema = [], {
                "index": s.index, "type": slide_type,
                "file": f"slide_{s.index:02d}_{slide_type}.svg",
                "texts": {}, "width_pt": round(page_w, 1), "height_pt": round(page_h, 1),
            }
            role_counts: dict[str, int] = {}
            shapes_by_id = {sh.shape_id: sh
                            for sh in iter_leaf_shapes(slides[s.index].shapes)}
            for t in s.texts:
                shape = shapes_by_id.get(t.shape_id)
                if shape is None:
                    continue
                role_counts[t.role] = role_counts.get(t.role, 0) + 1
                n = role_counts[t.role]
                name = f"{t.role}" + (f"_{n}" if n > 1 else "")
                font_pt, color, bold, align = _shape_text_style(shape)
                font_pt = font_pt or (28 if t.role == "title" else 14)

                x = t.left_in * 72
                y = t.top_in * 72
                w = t.width_in * 72
                anchor_x = {"left": x + 2, "middle": x + w / 2, "end": x + w - 2}[
                    align if align in ("middle", "end") else "left"]
                weight = ' font-weight="bold"' if bold else ""
                lines = max(t.n_paragraphs, 1)
                for ln in range(lines):
                    ph = f"{{{{{name}}}}}" if lines == 1 else f"{{{{{name}.{ln + 1}}}}}"
                    line_y = y + font_pt * 1.05 + ln * font_pt * 1.35
                    overlays.append(
                        f'<text x="{anchor_x:.1f}" y="{line_y:.1f}" '
                        f'font-family="{html.escape(font_family)}" '
                        f'font-size="{font_pt:.1f}" fill="{color}"'
                        f'{weight} text-anchor="{align if align in ("middle", "end") else "start"}"'
                        f'>{ph}</text>')

                slide_schema["texts"][name] = {
                    "role": t.role, "max_chars": t.max_chars,
                    "paragraphs": t.n_paragraphs, "sample": t.current_text[:100],
                }

            svg = (
                f'<svg xmlns="http://www.w3.org/2000/svg" '
                f'xmlns:xlink="http://www.w3.org/1999/xlink" '
                f'viewBox="0 0 {page_w:.1f} {page_h:.1f}">'
                f'<image x="0" y="0" width="{page_w:.1f}" height="{page_h:.1f}" '
                f'xlink:href="data:image/png;base64,{b64}"/>'
                + "".join(overlays) + "</svg>"
            )
            (out_dir / slide_schema["file"]).write_text(svg, encoding="utf-8")
            schema["slides"].append(slide_schema)

        doc.close()

    (out_dir / "schema.json").write_text(
        json.dumps(schema, ensure_ascii=False, indent=2), encoding="utf-8")
    return schema
