"""Skill 1: parse a Canva-exported .pptx template into a fill-spec.

Canva exports slides as plain text boxes and pictures (not PowerPoint
placeholder shapes), so this walks every shape, records what text/images
exist, and classifies each text box by role (title / subtitle / body /
caption) using font size and position heuristics. The resulting spec is
JSON-serializable so it can travel through an API or an LLM prompt.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


@dataclass
class TextElement:
    shape_id: int
    name: str
    role: str               # title | subtitle | body | caption
    current_text: str
    max_chars: int          # soft budget derived from the template's sample text
    n_paragraphs: int       # bullet lists keep one line per paragraph
    font_size_pt: float | None
    left_in: float
    top_in: float
    width_in: float
    height_in: float


@dataclass
class ImageElement:
    shape_id: int
    name: str
    width_px: int
    height_px: int
    aspect_ratio: float     # width / height
    left_in: float
    top_in: float


@dataclass
class SlideSpec:
    index: int
    texts: list[TextElement] = field(default_factory=list)
    images: list[ImageElement] = field(default_factory=list)


@dataclass
class TemplateSpec:
    source_path: str
    slide_width_in: float
    slide_height_in: float
    slides: list[SlideSpec] = field(default_factory=list)

    def to_dict(self) -> dict:
        return asdict(self)

    def to_json(self, **kwargs) -> str:
        return json.dumps(self.to_dict(), ensure_ascii=False, **kwargs)


_EMU_PER_PX = 9525  # at 96 dpi


def iter_leaf_shapes(shapes):
    """Yield non-group shapes, recursing into groups (Canva groups freely)."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _on_canvas(shape, slide_w_emu: int, slide_h_emu: int) -> bool:
    """False when the shape sits entirely outside the slide. Canva parks its
    'Mẹo:'/tip helper bubbles off-canvas; they must not be treated as
    fillable content."""
    if shape.left is None or shape.top is None:
        return True
    right = shape.left + (shape.width or 0)
    bottom = shape.top + (shape.height or 0)
    return right > 0 and bottom > 0 and shape.left < slide_w_emu and shape.top < slide_h_emu


def _max_font_size_pt(text_frame) -> float | None:
    sizes = [
        run.font.size.pt
        for para in text_frame.paragraphs
        for run in para.runs
        if run.font.size is not None
    ]
    return max(sizes) if sizes else None


# Canva's instructional templates ship visible helper bubbles; their text
# starts with the UI language's "Tip:" prefix. Extend per locale as needed.
DEFAULT_EXCLUDE_PREFIXES = ("Mẹo:", "Tip:", "Tips:", "Consejo:", "Astuce :")


def _is_navigation(text_frame) -> bool:
    """True when every run is a hyperlink — e.g. Canva's 'back to table of
    contents' buttons. Rewriting those breaks deck navigation."""
    runs = [r for p in text_frame.paragraphs for r in p.runs if r.text.strip()]
    return bool(runs) and all(r.hyperlink.address for r in runs)


def _classify_role(font_pt: float | None, top_in: float, slide_height_in: float,
                   text: str) -> str:
    """Heuristic role from font size, vertical position, and text length."""
    if font_pt is not None:
        if font_pt >= 28:
            return "title"
        if font_pt >= 18:
            return "subtitle"
        if font_pt <= 11:
            return "caption"
        return "body"
    # No explicit size (inherited from theme): fall back to position/length.
    if top_in < slide_height_in * 0.25 and len(text) <= 80:
        return "title"
    if len(text) <= 60:
        return "caption"
    return "body"


def parse_template(
    pptx_path: Union[str, Path],
    *,
    exclude_navigation: bool = True,
    exclude_prefixes: tuple[str, ...] = DEFAULT_EXCLUDE_PREFIXES,
) -> TemplateSpec:
    """Inspect a .pptx and return a TemplateSpec describing every fillable
    text box and picture, slide by slide.

    Skipped (left untouched by the filler): off-canvas shapes, fully
    hyperlinked text (navigation buttons) when exclude_navigation, and text
    starting with a Canva tip prefix from exclude_prefixes."""
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    slide_w_in = Emu(prs.slide_width).inches
    slide_h_in = Emu(prs.slide_height).inches
    spec = TemplateSpec(
        source_path=str(pptx_path),
        slide_width_in=round(slide_w_in, 2),
        slide_height_in=round(slide_h_in, 2),
    )

    for idx, slide in enumerate(prs.slides):
        slide_spec = SlideSpec(index=idx)
        # Off-canvas filtering happens at the top level only: nested shapes
        # report group-local coordinates that can't be compared to the slide.
        on_canvas = (
            s for s in slide.shapes
            if _on_canvas(s, prs.slide_width, prs.slide_height)
        )
        for shape in iter_leaf_shapes(on_canvas):
            if shape.has_text_frame and shape.text_frame.text.strip():
                tf = shape.text_frame
                text = tf.text
                if text.strip().startswith(exclude_prefixes):
                    continue
                if exclude_navigation and _is_navigation(tf):
                    continue
                font_pt = _max_font_size_pt(tf)
                top_in = Emu(shape.top).inches if shape.top is not None else 0.0
                slide_spec.texts.append(TextElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    role=_classify_role(font_pt, top_in, slide_h_in, text),
                    current_text=text,
                    max_chars=max(len(text), 20),
                    n_paragraphs=len([p for p in tf.paragraphs if p.text.strip()]),
                    font_size_pt=font_pt,
                    left_in=round(Emu(shape.left or 0).inches, 2),
                    top_in=round(top_in, 2),
                    width_in=round(Emu(shape.width or 0).inches, 2),
                    height_in=round(Emu(shape.height or 0).inches, 2),
                ))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                w_px = int((shape.width or 0) / _EMU_PER_PX)
                h_px = int((shape.height or 0) / _EMU_PER_PX)
                if w_px == 0 or h_px == 0:
                    continue
                slide_spec.images.append(ImageElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    width_px=w_px,
                    height_px=h_px,
                    aspect_ratio=round(w_px / h_px, 3),
                    left_in=round(Emu(shape.left or 0).inches, 2),
                    top_in=round(Emu(shape.top or 0).inches, 2),
                ))
        spec.slides.append(slide_spec)

    return spec
