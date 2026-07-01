"""Skill 19: AI-assisted template extraction — a .pptx design into a
category-style SVG template folder the dynamic flow can map to.

    extract_template_smart("deck.pptx", "purpose")
        -> templates/purpose_template/<CATEGORY>/{standard.svg,
                                                  category.json,
                                                  standard.schema.json}

Per slide:
  - render the slide's design to a background (LibreOffice), with CONTENT
    text blanked and CONTENT images removed — so decoration/background stays
    baked in, content areas are empty.
  - overlay {{placeholders}} for the content text and {{image}} slots for the
    content images.
  - images are classified content-vs-decoration by a size heuristic (large /
    full-bleed = decoration, kept; smaller framed = content slot).
  - one GPT-4o call writes the category name + per-field type/description, so
    the downstream mapper knows what each slot is for.

Needs LibreOffice (pptx -> PDF render) and PyMuPDF (PDF -> PNG).
"""

from __future__ import annotations

import base64
import html
import json
import re
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import get_client, TEXT_MODEL, collections_dir
from .svg_template_maker import _pptx_to_pdf, _shape_text_style
from .template_parser import iter_leaf_shapes, parse_template
from .usage import tracker

# an image bigger than this fraction of the slide is treated as decoration/
# background (kept in the render), not a fillable content slot.
_DECORATION_AREA = 0.30


def _classify_image(shape, slide_w_emu: int, slide_h_emu: int) -> str:
    """'decoration' (keep, baked into background) or 'content' (fillable slot)."""
    w = (shape.width or 0) / slide_w_emu if slide_w_emu else 0
    h = (shape.height or 0) / slide_h_emu if slide_h_emu else 0
    if w * h >= _DECORATION_AREA or w >= 0.9 or h >= 0.9:
        return "decoration"
    return "content"


_CLASSIFY_SYSTEM = """\
You classify the slides of a presentation template so a generator can reuse
them. For EVERY slide, return:
- category: a short UPPER_SNAKE_CASE name for the slide's function (e.g.
  TITLE_SLIDE, AGENDA, SECTION_HEADER, BULLETS, KPI_STATS, MEDIA_TEXT, QUOTE,
  TIMELINE, COMPARISON, SUMMARY, CTA, CONTACT). Reuse the same name for slides
  with the same function; suffix _2 if a distinct second design of the same kind.
- purpose: one sentence — when to use this slide.
- fields: for each text placeholder name given, its {type, desc}. type is one
  of title/subtitle/text/number/caption. desc says what the field is for
  (e.g. "the main statistic, a number like 78%", "a supporting bullet point").
Return ONLY JSON:
{"slides":[{"index":0,"category":"...","purpose":"...",
  "fields":{"<name>":{"type":"...","desc":"..."}}}]}
Cover every index.
"""


def _classify_slides(slide_infos: list) -> dict:
    """One GPT-4o call -> {index: {category, purpose, fields}}."""
    client = get_client()
    response = client.chat.completions.create(
        model=TEXT_MODEL,
        temperature=0.2,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": _CLASSIFY_SYSTEM},
            {"role": "user", "content": json.dumps(slide_infos, ensure_ascii=False)},
        ],
    )
    tracker.record_chat(response.usage)
    payload = json.loads(response.choices[0].message.content)
    return {s.get("index"): s for s in payload.get("slides", [])}


def _placeholder_name(role: str, counts: dict) -> str:
    counts[role] = counts.get(role, 0) + 1
    n = counts[role]
    return f"{role}" + (f"_{n}" if n > 1 else "")


def extract_template_smart(
    pptx_path: Union[str, Path],
    name: str,
    *,
    library_dir: Union[str, Path, None] = None,
    use_ai: bool = True,
) -> dict:
    """Convert a .pptx into templates/<name>_template/<CATEGORY>/ entries
    (standard.svg + category.json + standard.schema.json). Returns a summary."""
    import fitz

    pptx_path = Path(pptx_path)
    base = Path(library_dir) if library_dir is not None else collections_dir()
    out_root = base / f"{name}_template"
    out_root.mkdir(parents=True, exist_ok=True)

    spec = parse_template(pptx_path)
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    sw, sh = prs.slide_width, prs.slide_height
    page_w = sw / 914400 * 72
    page_h = sh / 914400 * 72

    # --- classify images per slide (content vs decoration) + plan placeholders
    fillable_text_ids = {s.index: {t.shape_id for t in s.texts} for s in spec.slides}
    per_slide = {}   # index -> {content_img_ids, text_names, img_names, text_meta}
    for s in spec.slides:
        content_imgs, img_names = [], []
        shapes_by_id = {sh_.shape_id: sh_ for sh_ in iter_leaf_shapes(slides[s.index].shapes)}
        for k, img in enumerate(s.images, 1):
            shp = shapes_by_id.get(img.shape_id)
            if shp is None:
                continue
            if _classify_image(shp, sw, sh) == "content":
                nm = f"image" + (f"_{len(img_names)+1}" if img_names else "")
                content_imgs.append((img.shape_id, nm, img))
                img_names.append(nm)
        per_slide[s.index] = {"content_imgs": content_imgs, "img_names": img_names}

    # --- AI classification (category + field descriptions)
    ai = {}
    if use_ai:
        infos = []
        for s in spec.slides:
            counts = {}
            fields = [{"name": _placeholder_name(t.role, counts),
                       "role": t.role, "sample": t.current_text[:60]} for t in s.texts]
            infos.append({"index": s.index, "text_fields": fields,
                          "content_images": len(per_slide[s.index]["img_names"])})
        try:
            ai = _classify_slides(infos)
        except Exception:
            ai = {}

    # --- build a render copy: blank content text, remove content images
    work = out_root / ".work"
    work.mkdir(exist_ok=True)
    blanked = work / f"blank_{pptx_path.stem}.pptx"
    pdf = work / f"blank_{pptx_path.stem}.pdf"
    from .slide_filler import set_text_preserving_format
    for idx, slide in enumerate(slides):
        content_ids = {sid for sid, _, _ in per_slide.get(idx, {}).get("content_imgs", [])}
        for shape in list(iter_leaf_shapes(slide.shapes)):
            if shape.shape_id in fillable_text_ids.get(idx, set()) and shape.has_text_frame:
                set_text_preserving_format(shape.text_frame, "")
            elif shape.shape_id in content_ids:
                shape._element.getparent().remove(shape._element)
    prs.save(str(blanked))
    _pptx_to_pdf(blanked, pdf)
    doc = fitz.open(str(pdf))

    # --- emit one category folder per slide
    seen_cats, summary = {}, []
    fresh = list(Presentation(str(pptx_path)).slides)  # original (for positions)
    for s in spec.slides:
        meta = ai.get(s.index, {})
        cat = re.sub(r"[^A-Z0-9_]", "_", str(meta.get("category", f"SLIDE_{s.index}")).upper())
        seen_cats[cat] = seen_cats.get(cat, 0) + 1
        if seen_cats[cat] > 1:
            cat = f"{cat}_{seen_cats[cat]}"
        cat_dir = out_root / cat
        cat_dir.mkdir(parents=True, exist_ok=True)

        # background PNG (decoration kept, content blanked)
        pix = doc[s.index].get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
        b64 = base64.b64encode(pix.tobytes("png")).decode()

        overlays, schema_fields = [], {}
        ai_fields = meta.get("fields", {})
        counts = {}
        shapes_by_id = {sh_.shape_id: sh_ for sh_ in iter_leaf_shapes(fresh[s.index].shapes)}
        for t in s.texts:
            nm = _placeholder_name(t.role, counts)
            shp = shapes_by_id.get(t.shape_id)
            font_pt, color, bold, align = _shape_text_style(shp) if shp else (None, "#000", False, "left")
            font_pt = font_pt or (28 if t.role == "title" else 14)
            x, y, w = t.left_in * 72, t.top_in * 72, t.width_in * 72
            anchor = {"left": x + 2, "middle": x + w / 2, "end": x + w - 2}[
                align if align in ("middle", "end") else "left"]
            marker = "\n".join(f"{{{{{nm}.{i+1}}}}}" for i in range(t.n_paragraphs)) \
                if t.n_paragraphs > 1 else f"{{{{{nm}}}}}"
            weight = ' font-weight="bold"' if bold else ""
            anchor_attr = align if align in ("middle", "end") else "start"
            font_family = html.escape(_FONT)
            box_w = max(w - 4, 12)
            for ln, part in enumerate(marker.split("\n")):
                ty = y + font_pt * 1.05 + ln * font_pt * 1.35
                overlays.append(
                    f'<text x="{anchor:.1f}" y="{ty:.1f}" data-w="{box_w:.0f}" '
                    f'font-family="{font_family}" font-size="{font_pt:.1f}" '
                    f'fill="{color}"{weight} '
                    f'text-anchor="{anchor_attr}">{part}</text>')
            fmeta = ai_fields.get(nm, {})
            schema_fields[nm] = {"type": fmeta.get("type", t.role),
                                 "desc": fmeta.get("desc", f"{t.role} text"),
                                 "max_chars": t.max_chars}
            if t.n_paragraphs > 1:
                schema_fields[nm]["lines"] = t.n_paragraphs

        # content image slots (as <image href="{{...}}">)
        for sid, nm, img in per_slide.get(s.index, {}).get("content_imgs", []):
            x, y = img.left_in * 72, img.top_in * 72
            w, h = img.width_px / 96 * 72, img.height_px / 96 * 72
            overlays.append(
                f'<image x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                f'preserveAspectRatio="xMidYMid slice" href="{{{{{nm}}}}}"/>')

        svg = (f'<svg xmlns="http://www.w3.org/2000/svg" '
               f'xmlns:xlink="http://www.w3.org/1999/xlink" '
               f'viewBox="0 0 {page_w:.1f} {page_h:.1f}">'
               f'<image x="0" y="0" width="{page_w:.1f}" height="{page_h:.1f}" '
               f'href="data:image/png;base64,{b64}"/>'
               + "".join(overlays) + "</svg>")
        (cat_dir / "standard.svg").write_text(svg, encoding="utf-8")
        purpose = meta.get("purpose", f"Slide {s.index}")
        (cat_dir / "category.json").write_text(json.dumps(
            {"description": purpose,
             "variants": {"standard": purpose}},
            ensure_ascii=False, indent=2), encoding="utf-8")
        (cat_dir / "standard.schema.json").write_text(json.dumps(
            {"fields": schema_fields}, ensure_ascii=False, indent=2), encoding="utf-8")
        # a compact, self-describing entry so the return value alone tells an
        # agent WHAT each category is for and WHAT each slot expects.
        fields_desc = {fn: {"type": fm.get("type"), "desc": fm.get("desc"),
                            "max_chars": fm.get("max_chars")}
                       for fn, fm in schema_fields.items()}
        summary.append({
            "slide": s.index,
            "category": cat,
            "purpose": purpose,                 # what this slide type is for
            "text_slots": len(schema_fields),
            "image_slots": len(per_slide.get(s.index, {}).get("img_names", [])),
            "fields": fields_desc,              # per-slot type + description
        })

    doc.close()
    import shutil
    shutil.rmtree(work, ignore_errors=True)   # drop the render scratch dir

    # single catalog manifest: one file an agent can read to know the whole
    # template (every category, its purpose, and its fillable fields).
    catalog = [{"category": e["category"], "purpose": e["purpose"],
                "fields": e["fields"]} for e in summary]
    (out_root / "library.json").write_text(
        json.dumps({"template": out_root.name, "categories": catalog},
                   ensure_ascii=False, indent=2), encoding="utf-8")

    return {"template_dir": str(out_root),
            "manifest": str(out_root / "library.json"),
            "categories": catalog,       # descriptions the picker/agent needs
            "slides": summary,
            "usage": tracker.snapshot().report()}


_FONT = "Helvetica, Arial, sans-serif"
