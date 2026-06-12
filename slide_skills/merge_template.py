"""Skill 11: {{placeholder}} merge templates.

Converts a deck's fillable content into named {{placeholders}} plus a JSON
schema, so generation becomes a clean three-step contract:

    make_placeholder_template  -- slides' sample text -> {{s0.title}} markers,
                                  schema.json describing every placeholder
    generate_merge_data        -- GPT-4o returns {placeholder: value} JSON
                                  honoring the schema's budgets
    render_placeholders        -- pure code: substitute values (and images)
                                  into the template. Deterministic, offline.

The schema is what an API/chatbot stores and validates against; the AI never
touches the .pptx and the renderer never touches the AI.
"""

from __future__ import annotations

import io
import json
import re
from pathlib import Path
from typing import Union

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import get_client, TEXT_MODEL
from .slide_filler import replace_picture_image, set_text_preserving_format, shrink_to_fit
from .template_parser import iter_leaf_shapes, parse_template
from .usage import tracker

PLACEHOLDER_RE = re.compile(r"\{\{([^{}]+)\}\}")


def _placeholder_png(name: str, width_px: int, height_px: int) -> bytes:
    """A visible image placeholder: gray block, crossed diagonals, the
    {{name}} and pixel size printed on it."""
    scale = min(1.0, 800 / max(width_px, height_px, 1))
    w, h = max(int(width_px * scale), 96), max(int(height_px * scale), 96)
    img = Image.new("RGB", (w, h), (232, 236, 239))
    draw = ImageDraw.Draw(img)
    draw.line([(0, 0), (w, h)], fill=(200, 206, 212), width=2)
    draw.line([(0, h), (w, 0)], fill=(200, 206, 212), width=2)
    draw.rectangle([0, 0, w - 1, h - 1], outline=(150, 158, 166), width=3)
    label = f"{{{{{name}}}}}"
    size_label = f"{width_px}x{height_px}px"
    lw = draw.textlength(label)
    draw.rectangle([w / 2 - lw / 2 - 8, h / 2 - 22, w / 2 + lw / 2 + 8, h / 2 + 18],
                   fill=(255, 255, 255))
    draw.text((w / 2 - lw / 2, h / 2 - 14), label, fill=(60, 70, 80))
    draw.text((w / 2 - draw.textlength(size_label) / 2, h / 2 + 2),
              size_label, fill=(130, 140, 150))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _slide_prefix(index: int, types: list[str] | None) -> str:
    if types and index < len(types) and not types[index].startswith("_"):
        return types[index]
    return f"s{index}"


def make_placeholder_template(
    src: Union[str, Path],
    dst: Union[str, Path],
    *,
    types: list[str] | None = None,
    mark_images: bool = True,
) -> dict:
    """Rewrite every fillable text box as {{name}} markers (formatting kept)
    and produce the schema. Saves the schema next to dst as <dst>.schema.json
    and returns it.

    Multi-paragraph boxes become {{name.1}}..{{name.N}} — supply a list when
    rendering. With mark_images (default), each picture placeholder is
    replaced by a gray block showing its {{name}} and size — the template
    becomes fully explicit and sheds the original stock photos."""
    src, dst = Path(src), Path(dst)
    if types is None:
        manifest = src.with_suffix(".json")
        if manifest.exists():
            types = json.loads(manifest.read_text(encoding="utf-8")).get("types")

    spec = parse_template(src)
    prs = Presentation(str(src))
    schema: dict = {"text": {}, "images": {}}

    slides = list(prs.slides)
    for s in spec.slides:
        prefix = _slide_prefix(s.index, types)
        fillable = {t.shape_id: t for t in s.texts}
        role_counts: dict[str, int] = {}
        shapes_by_id = {sh.shape_id: sh for sh in iter_leaf_shapes(slides[s.index].shapes)}

        for shape_id, t in fillable.items():
            shape = shapes_by_id.get(shape_id)
            if shape is None or not shape.has_text_frame:
                continue
            role_counts[t.role] = role_counts.get(t.role, 0) + 1
            n = role_counts[t.role]
            name = f"{prefix}.{t.role}" + (f"_{n}" if n > 1 else "")

            if t.n_paragraphs > 1:
                marker = "\n".join(
                    f"{{{{{name}.{i + 1}}}}}" for i in range(t.n_paragraphs))
            else:
                marker = f"{{{{{name}}}}}"
            set_text_preserving_format(shape.text_frame, marker)

            schema["text"][name] = {
                "slide": s.index,
                "shape_id": shape_id,
                "role": t.role,
                "max_chars": t.max_chars,
                "paragraphs": t.n_paragraphs,
                "sample": t.current_text[:120],
            }

        for k, img in enumerate(s.images, 1):
            name = f"{prefix}.image" + (f"_{k}" if len(s.images) > 1 else "")
            schema["images"][name] = {
                "slide": s.index,
                "shape_id": img.shape_id,
                "aspect_ratio": img.aspect_ratio,
                "width_px": img.width_px,
                "height_px": img.height_px,
            }
            if mark_images:
                shape = shapes_by_id.get(img.shape_id)
                if shape is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    replace_picture_image(
                        shape, _placeholder_png(name, img.width_px, img.height_px))

    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    schema_path = dst.with_suffix(".schema.json")
    schema_path.write_text(
        json.dumps(schema, ensure_ascii=False, indent=2), encoding="utf-8")
    return schema


def load_schema(template_path: Union[str, Path]) -> dict:
    return json.loads(
        Path(template_path).with_suffix(".schema.json").read_text(encoding="utf-8"))


def render_placeholders(
    template_path: Union[str, Path],
    data: dict,
    output_path: Union[str, Path],
    *,
    images: dict[str, bytes] | None = None,
    fit_text: bool = True,
) -> list[str]:
    """Substitute {{placeholders}} with data values (lists fill .1/.2/...
    markers) and image bytes by placeholder name. Returns the names of
    placeholders that were left unfilled."""
    schema = load_schema(template_path)

    flat: dict[str, str] = {}
    for key, value in data.items():
        if isinstance(value, (list, tuple)):
            for i, item in enumerate(value, 1):
                flat[f"{key}.{i}"] = str(item)
        else:
            flat[key] = str(value)

    budgets = {name: meta["max_chars"] for name, meta in schema["text"].items()}
    image_targets = {
        name: (meta["slide"], meta["shape_id"])
        for name, meta in schema["images"].items()
    }
    images = images or {}

    prs = Presentation(str(template_path))
    slides = list(prs.slides)
    unfilled: set[str] = set()

    for slide in slides:
        for shape in iter_leaf_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            frame_names: set[str] = set()
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{{" not in run.text:
                        continue

                    def sub(match):
                        key = match.group(1).strip()
                        frame_names.add(key.rsplit(".", 1)[0] if key.rsplit(".", 1)[-1].isdigit() else key)
                        if key in flat:
                            return flat[key]
                        unfilled.add(key)
                        return ""

                    run.text = PLACEHOLDER_RE.sub(sub, run.text)
            if fit_text:
                for name in frame_names:
                    budget = budgets.get(name)
                    text_len = len(shape.text_frame.text)
                    if budget and text_len > budget:
                        shrink_to_fit(shape.text_frame, text_len, budget)

    by_id = [
        {sh.shape_id: sh for sh in iter_leaf_shapes(sl.shapes)} for sl in slides
    ]
    for name, png in images.items():
        if name in image_targets:
            idx, shape_id = image_targets[name]
            shape = by_id[idx].get(shape_id)
            if shape is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                replace_picture_image(shape, png)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return sorted(unfilled)


def generate_merge_data(
    template_path: Union[str, Path],
    brief: str,
    *,
    language: str | None = None,
    temperature: float = 0.7,
) -> tuple[dict, dict]:
    """GPT-4o fills the schema: returns (text_data, image_prompts).
    text_data values are strings, or lists for multi-paragraph placeholders.
    image_prompts maps image placeholder names to DALL-E prompts."""
    schema = load_schema(template_path)

    text_spec = {
        name: {"role": m["role"], "max_chars": m["max_chars"],
               "paragraphs": m["paragraphs"], "sample": m["sample"]}
        for name, m in schema["text"].items()
    }
    image_spec = {
        name: {"aspect_ratio": m["aspect_ratio"]}
        for name, m in schema["images"].items()
    }

    system = (
        "You fill presentation templates. You get a brief and a placeholder "
        "schema. Return a value for EVERY placeholder.\n"
        "Rules:\n"
        "- Stay within each placeholder's max_chars (hard limit).\n"
        "- paragraphs > 1 means return an ARRAY with exactly that many short "
        "strings; otherwise return a single string.\n"
        "- Match the role and the sample's tone/shape; write in the brief's "
        "language unless told otherwise.\n"
        "- For every image placeholder return a concrete DALL-E scene prompt "
        "(subject, style, mood, colors), no text in images.\n"
        'Return ONLY JSON: {"text": {"<name>": "..." | ["...", "..."]}, '
        '"images": {"<name>": "prompt"}}'
    )
    user = (f"Brief:\n{brief}\n\nText placeholders:\n"
            f"{json.dumps(text_spec, ensure_ascii=False)}\n\n"
            f"Image placeholders:\n{json.dumps(image_spec, ensure_ascii=False)}")
    if language:
        user += f"\n\nWrite all content in {language}."

    client = get_client()
    response = client.chat.completions.create(
        model=TEXT_MODEL,
        temperature=temperature,
        response_format={"type": "json_object"},
        messages=[{"role": "system", "content": system},
                  {"role": "user", "content": user}],
    )
    tracker.record_chat(response.usage)
    payload = json.loads(response.choices[0].message.content)

    text_data = {
        k: v for k, v in (payload.get("text") or {}).items()
        if k in schema["text"]
    }
    image_prompts = {
        k: str(v) for k, v in (payload.get("images") or {}).items()
        if k in schema["images"]
    }
    return text_data, image_prompts
