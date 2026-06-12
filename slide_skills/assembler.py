"""Skill 8: assemble a deck from a slide library.

The library is one .pptx whose slides are reusable designs in a known order
(see DEFAULT_LIBRARY_TYPES). Given a planned sequence of types — repeats
allowed — this picks, duplicates, and reorders library slides into a new
deck, dropping unused ones. Duplication clones the slide XML plus its image
relationships, so repeated slides are independently fillable.

A <library>.json manifest next to the .pptx overrides the type order:
    {"types": ["title", "agenda", ...],          # one entry per slide, in order
     "descriptions": {"agenda": "what it's for"}} # optional, guides the planner

Type names starting with "_" mark slides the planner must never use
(e.g. Canva tutorial pages kept in the file).
"""

from __future__ import annotations

import copy
import json
from pathlib import Path
from typing import Union

from pptx import Presentation

DEFAULT_LIBRARY_TYPES = [
    "title", "agenda", "section", "concept", "bullets",
    "statistic", "graph", "comparison", "quote", "summary",
]

_RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_SKIP_RELTYPES = ("slideLayout", "notesSlide")


def load_library_meta(library_path: Union[str, Path]) -> tuple[list[str], dict[str, str]]:
    """(type per slide, {type: description}) from the manifest if present,
    otherwise DEFAULT_LIBRARY_TYPES trimmed to the slide count."""
    library_path = Path(library_path)
    manifest = library_path.with_suffix(".json")
    n_slides = len(Presentation(str(library_path)).slides._sldIdLst)

    if manifest.exists():
        data = json.loads(manifest.read_text(encoding="utf-8"))
        types = [str(t) for t in data["types"]]
        if len(types) != n_slides:
            raise ValueError(
                f"{manifest.name} lists {len(types)} types but the library "
                f"has {n_slides} slides."
            )
        descriptions = {str(k): str(v)
                        for k, v in (data.get("descriptions") or {}).items()}
        return types, descriptions

    if n_slides > len(DEFAULT_LIBRARY_TYPES):
        raise ValueError(
            f"Library has {n_slides} slides but no manifest; the default "
            f"order only covers {len(DEFAULT_LIBRARY_TYPES)}. Add a "
            f"{manifest.name} manifest."
        )
    return DEFAULT_LIBRARY_TYPES[:n_slides], {}


def load_library_types(library_path: Union[str, Path]) -> list[str]:
    return load_library_meta(library_path)[0]


def _strip_slide_links(slide) -> None:
    """Remove slide-to-slide hyperlink relationships from a slide.

    Canva decks link agenda items and 'back' buttons to other slides. After
    assembly reorders/prunes slides those links point at the wrong place, and
    the dangling cross-references corrupt the save (duplicate zip entries).
    The r:id is removed but the hlinkClick element stays, so the template
    parser still recognizes the text as navigation and leaves it alone."""
    part = slide.part
    slide_rids = {
        rel.rId for rel in part.rels.values()
        if not rel.is_external and rel.reltype.endswith("/slide")
    }
    if not slide_rids:
        return
    for node in slide._element.iter():
        if node.tag.endswith("}hlinkClick") or node.tag.endswith("}hlinkHover"):
            rid_attr = "{%s}id" % _RELS_NS
            if node.get(rid_attr) in slide_rids:
                del node.attrib[rid_attr]
    for rid in slide_rids:
        part.drop_rel(rid)


def _duplicate_slide(prs, source_slide):
    """Clone a slide (shapes + image/media relationships) onto a new slide
    appended at the end. Returns the new slide."""
    dest = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    rid_map = {}
    for rel in list(source_slide.part.rels.values()):
        if any(rel.reltype.endswith(s) for s in _SKIP_RELTYPES):
            continue
        if rel.is_external:
            rid_map[rel.rId] = dest.part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref)
        else:
            rid_map[rel.rId] = dest.part.relate_to(rel.target_part, rel.reltype)

    for shape in source_slide.shapes:
        element = copy.deepcopy(shape._element)
        for node in element.iter():
            for attr, value in list(node.attrib.items()):
                if attr.startswith("{%s}" % _RELS_NS) and value in rid_map:
                    node.attrib[attr] = rid_map[value]
        dest.shapes._spTree.append(element)
    return dest


def build_deck_from_library(
    library_path: Union[str, Path],
    type_sequence: list[str],
    output_path: Union[str, Path],
    *,
    library_types: list[str] | None = None,
) -> Path:
    """Create output_path containing one slide per entry in type_sequence,
    in that order, cloned from the library."""
    library_types = library_types or load_library_types(library_path)
    prs = Presentation(str(library_path))
    originals = list(prs.slides)
    if len(library_types) != len(originals):
        raise ValueError(
            f"library_types has {len(library_types)} entries for "
            f"{len(originals)} slides.")
    by_type = {t: s for t, s in zip(library_types, originals)}

    unknown = [t for t in type_sequence if t not in by_type]
    if unknown:
        raise ValueError(f"Types not in library: {unknown}. "
                         f"Available: {sorted(by_type)}")

    # First use takes the original slide; repeats get clones (appended last,
    # which is fine — ordering is fixed below).
    ordered, used = [], set()
    for t in type_sequence:
        if t in used:
            ordered.append(_duplicate_slide(prs, by_type[t]))
        else:
            ordered.append(by_type[t])
            used.add(t)

    # Pair each sldId entry with its slide, drop unused slides, then move the
    # kept entries into planned order (lxml append relocates the node).
    sld_id_lst = prs.slides._sldIdLst
    pairs = list(zip(list(sld_id_lst), list(prs.slides)))
    keep = {id(s) for s in ordered}
    sld_id_by_slide = {}
    for sld_id, slide in pairs:
        if id(slide) in keep:
            sld_id_by_slide[id(slide)] = sld_id
        else:
            r_id = sld_id.rId
            sld_id_lst.remove(sld_id)
            prs.part.drop_rel(r_id)
    for slide in ordered:
        sld_id_lst.append(sld_id_by_slide[id(slide)])
        _strip_slide_links(slide)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
