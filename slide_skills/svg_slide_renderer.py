"""Skill 15: render filled SVG slide templates into a deck.

Counterpart of svg_template_maker: takes the folder of {{placeholder}} SVGs
plus per-slide data, rasterizes each filled SVG, and assembles a .pptx of
full-bleed slide images (or just returns the PNGs).

    render_svg_slide(svg_path, {"title": "...", "body": ["a", "b"]}) -> PNG
    render_svg_deck("svg_templates/untitled", deck_data, "out/deck.pptx")

deck_data is a list of {"type": <slide type>, "texts": {...}} entries —
slides may repeat types and appear in any order.
"""

from __future__ import annotations

import io
import json
import re
from pathlib import Path
from typing import Union

import resvg_py
from pptx import Presentation
from pptx.util import Emu

_PLACEHOLDER_RE = re.compile(r"\{\{([^{}]+)\}\}")


def _fill(svg: str, data: dict) -> str:
    flat: dict[str, str] = {}
    for key, value in (data or {}).items():
        if isinstance(value, (list, tuple)):
            for i, item in enumerate(value, 1):
                flat[f"{key}.{i}"] = str(item)
        else:
            flat[key] = str(value)

    def sub(match):
        return flat.get(match.group(1).strip(), "")

    return _PLACEHOLDER_RE.sub(sub, svg)


def render_svg_slide(
    svg_path: Union[str, Path],
    data: dict,
    *,
    width_px: int = 1920,
) -> bytes:
    svg = Path(svg_path).read_text(encoding="utf-8")
    return bytes(resvg_py.svg_to_bytes(svg_string=_fill(svg, data),
                                       width=width_px))


def render_svg_deck(
    template_dir: Union[str, Path],
    deck_data: list[dict],
    output_path: Union[str, Path],
    *,
    width_px: int = 1920,
) -> Path:
    template_dir = Path(template_dir)
    schema = json.loads((template_dir / "schema.json").read_text(encoding="utf-8"))
    by_type = {s["type"]: s for s in schema["slides"]}

    first = schema["slides"][0]
    aspect = first["width_pt"] / first["height_pt"]

    prs = Presentation()
    prs.slide_width = Emu(int(12192000))                 # 13.33 in
    prs.slide_height = Emu(int(12192000 / aspect))
    blank = prs.slide_layouts[6]

    for entry in deck_data:
        slide_type = entry.get("type")
        meta = by_type.get(slide_type)
        if meta is None:
            raise ValueError(f"Unknown slide type {slide_type!r}. "
                             f"Available: {sorted(by_type)}")
        png = render_svg_slide(template_dir / meta["file"],
                               entry.get("texts", {}), width_px=width_px)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(png), 0, 0,
                                 prs.slide_width, prs.slide_height)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
