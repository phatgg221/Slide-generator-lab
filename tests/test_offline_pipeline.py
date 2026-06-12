"""Offline test of the pptx mechanics — no OpenAI key needed.

Builds a fake "Canva template" (styled text boxes + a picture), then runs
parse_template -> (stubbed content) -> fill_template and asserts the text
landed, formatting survived, and the picture bitmap was swapped in place.

Run:  python tests/test_offline_pipeline.py
"""

import io
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slide_skills import parse_template, fill_template
from slide_skills.content_generator import GeneratedDeckContent, GeneratedSlideContent

TMP = Path(__file__).resolve().parent / "tmp"
TMP.mkdir(exist_ok=True)


def make_png(color, size=(640, 360)) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, format="PNG")
    return buf.getvalue()


def build_fake_template(path: Path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Sample Title Goes Here"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)

    bullets = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5), Inches(3))
    tf = bullets.text_frame
    first = tf.paragraphs[0].add_run()
    first.text = "First sample bullet point here"
    first.font.size = Pt(16)
    for text in ["Second sample bullet point", "Third sample bullet point"]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(16)

    img_path = TMP / "placeholder.png"
    img_path.write_bytes(make_png((200, 200, 200)))
    slide.shapes.add_picture(str(img_path), Inches(6), Inches(2), Inches(3.5), Inches(2.6))

    # Canva-style grouped content: must be found inside the group
    group = slide.shapes.add_group_shape()
    grouped = group.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4), Inches(0.6))
    grouped_run = grouped.text_frame.paragraphs[0].add_run()
    grouped_run.text = "Grouped caption text"
    grouped_run.font.size = Pt(10)

    # Canva-style off-canvas tip bubble: must be ignored
    tip = slide.shapes.add_textbox(Inches(-6), Inches(1), Inches(4), Inches(1))
    tip.text_frame.paragraphs[0].add_run().text = "Tip: drag your image here"

    prs.save(str(path))


def main():
    template = TMP / "fake_template.pptx"
    output = TMP / "filled.pptx"
    build_fake_template(template)

    # --- parse ---
    spec = parse_template(template)
    assert len(spec.slides) == 1, "expected 1 slide"
    s = spec.slides[0]
    assert len(s.texts) == 3, f"expected 3 text boxes, got {len(s.texts)}"
    assert len(s.images) == 1, f"expected 1 picture, got {len(s.images)}"
    all_text = [t.current_text for t in s.texts]
    assert "Grouped caption text" in all_text, "grouped text box not found"
    assert not any("Tip:" in t for t in all_text), "off-canvas tip was not filtered"

    title_el = next(t for t in s.texts if t.role == "title")
    body_el = next(t for t in s.texts if t.n_paragraphs == 3)
    assert body_el.n_paragraphs == 3, f"expected 3 paragraphs, got {body_el.n_paragraphs}"
    pic_el = s.images[0]
    assert 1.3 < pic_el.aspect_ratio < 1.4, f"unexpected aspect {pic_el.aspect_ratio}"
    print(f"✓ parse: title={title_el.shape_id}, bullets={body_el.shape_id} "
          f"({body_el.n_paragraphs} paras), picture={pic_el.shape_id} "
          f"(aspect {pic_el.aspect_ratio})")

    # --- fill with stubbed "generated" content (4 lines: clones a paragraph) ---
    content = GeneratedDeckContent(slides=[GeneratedSlideContent(
        index=0,
        texts={
            title_el.shape_id: "AI Farming Drones",
            body_el.shape_id: "Cut water use by 30%\nSpot crop disease early\nMap fields in minutes\nWorks offline in the field",
        },
        images={pic_el.shape_id: "unused prompt"},
    )])
    new_png = make_png((30, 120, 60))
    fill_template(template, content, output, images={(0, pic_el.shape_id): new_png})

    # --- verify ---
    prs = Presentation(str(output))
    slide = prs.slides[0]
    by_id = {sh.shape_id: sh for sh in slide.shapes}

    title_tf = by_id[title_el.shape_id].text_frame
    assert title_tf.text == "AI Farming Drones", title_tf.text
    run = title_tf.paragraphs[0].runs[0]
    assert run.font.size == Pt(40) and run.font.bold, "title formatting lost"
    assert run.font.color.rgb == RGBColor(0x1E, 0x27, 0x61), "title color lost"

    body_tf = by_id[body_el.shape_id].text_frame
    lines = [p.text for p in body_tf.paragraphs]
    assert len(lines) == 4 and lines[3] == "Works offline in the field", lines
    assert body_tf.paragraphs[3].runs[0].font.size == Pt(16), "cloned bullet formatting lost"

    pic = by_id[pic_el.shape_id]
    assert pic.image.blob == new_png, "picture bitmap was not replaced"
    assert pic.width == Inches(3.5), "picture geometry changed"

    print("✓ fill: text replaced, formatting + cloned-bullet style preserved")
    print("✓ image: bitmap swapped, geometry untouched")
    print(f"\nAll offline checks passed. Output: {output}")


if __name__ == "__main__":
    main()
