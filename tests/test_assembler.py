"""Offline test for the library assembler — no API key needed.

Builds a 4-slide fake library (one text box + image per slide), assembles a
deck with reordering, repeats, and omissions, then verifies slide order,
clone independence, and image survival.

Run:  python tests/test_assembler.py
"""

import io
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from slide_skills.assembler import build_deck_from_library
from slide_skills.slide_filler import set_text_preserving_format

TMP = Path(__file__).resolve().parent / "tmp"
TMP.mkdir(exist_ok=True)

TYPES = ["title", "concept", "statistic", "summary"]
COLORS = {"title": (200, 0, 0), "concept": (0, 200, 0),
          "statistic": (0, 0, 200), "summary": (200, 200, 0)}


def build_fake_library(path: Path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for t in TYPES:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"SLIDE TYPE {t}"
        run.font.size = Pt(32)
        img = TMP / f"lib_{t}.png"
        buf = io.BytesIO()
        Image.new("RGB", (320, 200), COLORS[t]).save(buf, format="PNG")
        img.write_bytes(buf.getvalue())
        slide.shapes.add_picture(str(img), Inches(1), Inches(3), Inches(3), Inches(2))
    prs.save(str(path))


def slide_text(slide) -> str:
    return " ".join(sh.text_frame.text for sh in slide.shapes
                    if sh.has_text_frame and sh.text_frame.text.strip())


def main():
    library = TMP / "fake_library.pptx"
    output = TMP / "assembled.pptx"
    build_fake_library(library)

    # repeats 'concept' twice, omits 'statistic', reorders
    sequence = ["title", "concept", "summary", "concept"]
    build_deck_from_library(library, sequence, output, library_types=TYPES)

    prs = Presentation(str(output))
    slides = list(prs.slides)
    assert len(slides) == 4, f"expected 4 slides, got {len(slides)}"

    got = [slide_text(s) for s in slides]
    expected = [f"SLIDE TYPE {t}" for t in sequence]
    assert got == expected, f"\nexpected {expected}\ngot      {got}"
    print(f"✓ order: {[t for t in sequence]}")

    assert not any("statistic" in t for t in got), "omitted slide leaked in"
    print("✓ omitted type dropped")

    # both 'concept' slides carry the green image, independently
    pics = [[sh for sh in s.shapes if sh.shape_type == 13] for s in slides]
    assert all(len(p) == 1 for p in pics), "a slide lost its picture"
    c1, c2 = pics[1][0], pics[3][0]
    assert c1.image.blob == c2.image.blob, "clone image differs from original"
    print("✓ images survive cloning")

    # editing the clone must not edit the original
    set_text_preserving_format(
        [sh for sh in slides[3].shapes if sh.has_text_frame][0].text_frame,
        "EDITED CLONE")
    assert slide_text(slides[1]) == "SLIDE TYPE concept", "edit leaked to original"
    assert slide_text(slides[3]) == "EDITED CLONE"
    prs.save(str(TMP / "assembled_edited.pptx"))
    print("✓ clones are independent")

    print(f"\nAll assembler checks passed. Output: {output}")


if __name__ == "__main__":
    main()
